from flask import Flask, request, send_file
from openai import OpenAI
from pptx import Presentation
import io
from flask_cors import CORS
import os
from dotenv import load_dotenv
load_dotenv()

api_key = os.getenv('OPENAI_API_KEY')
if not api_key:
    raise ValueError(
        "OpenAI API key is missing. Please set it in the .env file.")
app = Flask(__name__)

CORS(app)
client = OpenAI(api_key=api_key)



def generate_presentation_outline(topic, num_slides=5, layout="Varied"):
    prompt = f"Generate a PowerPoint outline on the topic '{topic}'. The presentation should have {num_slides} slides with a {layout} layout."

    response = client.chat.completions.create(
        model="gpt-4", 
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ]
    )
    return response.choices[0].message.content.strip().split("\n")


def create_presentation(outline, num_slides):
    prs = Presentation()

    slide_data = parse_outline(outline)
    slide_data = slide_data[1:]

    for slide_title, slide_content in slide_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.shapes.placeholders[1]

        title.text = slide_title
        add_detailed_content(content, slide_content)

    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

def parse_outline(outline):
    slide_data = []
    current_slide_title = None
    current_slide_content = []

    for line in outline:
        line = line.strip()

        if line.startswith("Slide"):
            if current_slide_title:
                slide_data.append(
                    (current_slide_title, "\n".join(current_slide_content)))

            current_slide_title = line.split(":")[1].strip()
            current_slide_content = []
        elif line:
            current_slide_content.append(line)

    if current_slide_title:
        slide_data.append(
            (current_slide_title, "\n".join(current_slide_content)))

    return slide_data



def add_detailed_content(content_placeholder, content_text):
    content_lines = content_text.split("\n")
    for line in content_lines:
        if line.strip():
            p = content_placeholder.text_frame.add_paragraph()
            p.text = line.strip()


@app.route('/generate', methods=['POST'])
def generate_ppt():
    data = request.get_json()
    topic = data.get('topic')
    num_slides = int(data.get('numSlides', 5))
    layout = data.get('layout', "Varied")

    outline = generate_presentation_outline(topic, num_slides, layout)
    print(outline)
    pptx_io = create_presentation(outline, num_slides)

    return send_file(pptx_io, as_attachment=True, download_name="generated_presentation.pptx", mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")


if __name__ == '__main__':
    app.run(host='127.0.0.1', port=5001)
